import logging
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

logger = logging.getLogger(__name__)

def extract_pptx(temp_path: str) -> str:
    if Presentation is None:
        logger.error("python-pptx no está instalado. No se puede procesar PPTX.")
        return ""
    
    try:
        prs = Presentation(temp_path)
        full_text = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
            
            if slide_text:
                full_text.append(f"--- Diapositiva {i+1} ---\n" + "\n".join(slide_text))
        
        return "\n\n".join(full_text)
    except Exception as e:
        logger.error(f"Error extrayendo texto de PPTX: {e}")
        return ""
